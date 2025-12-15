"""Export helpers for InfluenceOpt"""
import csv
from typing import Iterable


def export_selection_csv(path: str, selected_ids: Iterable[str]) -> None:
    with open(path, 'w', newline='', encoding='utf-8') as f:
        w = csv.writer(f)
        w.writerow(['id'])
        for i in selected_ids:
            w.writerow([i])


def export_brief_pptx(path: str, scenario, screenshot_path: str | None = None) -> None:
    """Create a minimal PPTX brief with one slide containing screenshot and key stats."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        raise RuntimeError("python-pptx is required to export PPTX briefs")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    left = Inches(0.5)
    top = Inches(0.5)
    if screenshot_path:
        slide.shapes.add_picture(screenshot_path, left, top, width=Inches(6))
        top = Inches(4.0)
    tx = slide.shapes.add_textbox(left, top, width=Inches(6), height=Inches(2))
    tf = tx.text_frame
    tf.text = f"Scenario: {scenario.name}"
    tf.add_paragraph().text = f"Budget: {scenario.params.get('budget', '')}"
    tf.add_paragraph().text = f"Reach: {scenario.result.get('reach', '')}"
    tf.add_paragraph().text = f"Risk: {scenario.params.get('risk_max', '')}"
    prs.save(path)
